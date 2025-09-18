from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os
import re
import textwrap

# --- Slide layout ---
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(7.5)
TITLE_TOP = Inches(0.1)
TITLE_HEIGHT = Inches(0.9)
CONTENT_TOP = TITLE_TOP + TITLE_HEIGHT + Inches(0.2)
CONTENT_LEFT = Inches(0.5)
CONTENT_WIDTH = SLIDE_WIDTH - 2 * CONTENT_LEFT
CONTENT_HEIGHT = SLIDE_HEIGHT - CONTENT_TOP - Inches(0.5)

# --- Fonts and Colors ---
TITLE_FONT_SIZE = Pt(30)
BULLET_FONT_SIZE = Pt(16)
TABLE_FONT_SIZE = Pt(13)
BULLET_INDENT = Inches(0.25)
FOOTER_FONT_SIZE = Pt(11)
FOOTER_COLOR = RGBColor(228, 0, 43)

TITLE_BG_COLOR = RGBColor(0, 32, 91)
TITLE_FONT_COLOR = RGBColor(228, 0, 43)
CONTENT_BG_COLOR = RGBColor(247, 248, 250)
TABLE_HEADER_BG = RGBColor(228, 0, 43)
TABLE_HEADER_FONT_COLOR = RGBColor(255, 255, 255)
TABLE_ROW_ALT_BG = RGBColor(240, 240, 240)
TEXT_COLOR = RGBColor(30, 30, 30)

def flatten_to_string(item):
    """Recursively flattens tuples/lists and returns a string."""
    if isinstance(item, (tuple, list)):
        return " ".join(flatten_to_string(i) for i in item)
    return str(item)

def add_footer(slide, page_num):
    # Add left aligned "HPGPT"
    left_box = slide.shapes.add_textbox(Inches(0.5), SLIDE_HEIGHT - Inches(0.30), Inches(3), Pt(FOOTER_FONT_SIZE.pt + 8))
    left_tf = left_box.text_frame
    left_tf.clear()
    p_left = left_tf.paragraphs[0]
    p_left.text = "HPGPT"
    p_left.font.size = FOOTER_FONT_SIZE
    p_left.font.bold = True
    p_left.font.name = "Segoe UI"
    p_left.font.color.rgb = FOOTER_COLOR
    p_left.alignment = PP_ALIGN.LEFT

    # Add right aligned page number
    right_box_width = Inches(3)
    right_box_left = SLIDE_WIDTH - Inches(0.5) - right_box_width
    right_box = slide.shapes.add_textbox(right_box_left, SLIDE_HEIGHT - Inches(0.30), right_box_width, Pt(FOOTER_FONT_SIZE.pt + 8))
    right_tf = right_box.text_frame
    right_tf.clear()
    p_right = right_tf.paragraphs[0]
    p_right.text = f"Page {page_num}"
    p_right.font.size = FOOTER_FONT_SIZE
    p_right.font.bold = True
    p_right.font.name = "Segoe UI"
    p_right.font.color.rgb = FOOTER_COLOR
    p_right.alignment = PP_ALIGN.RIGHT

def add_cover_slide(prs, filename_title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_fill = slide.background.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = RGBColor(240, 240, 255)
    textbox = slide.shapes.add_textbox(Inches(1.5), SLIDE_HEIGHT / 2 - TITLE_HEIGHT / 2,
                                       SLIDE_WIDTH - Inches(3), TITLE_HEIGHT)

    from textwrap import fill

    def format_title_properly(title: str) -> str:
        title = title.strip()
        title = re.sub(r'\s*:\s*', ': ', title)
        title = re.sub(r'\s*-\s*', ' - ', title)
        title_words = title.split()

        def cap_word(i, w):
            lw = w.lower()
            if i == 0 or i == len(title_words) - 1:
                return w.capitalize()
            if lw in {'and', 'or', 'the', 'of', 'in', 'on', 'with', 'a', 'an', 'to', 'for'}:
                return lw
            return w.capitalize()

        title_cased = ' '.join(cap_word(i, w) for i, w in enumerate(title_words))
        return title_cased

    raw_title = flatten_to_string(filename_title)
    formatted_title = format_title_properly(raw_title)
    wrapped_title = fill(formatted_title, width=40)

    tf = textbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = wrapped_title
    p.font.size = TITLE_FONT_SIZE
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.font.color.rgb = TITLE_FONT_COLOR
    p.alignment = PP_ALIGN.CENTER
    p.word_wrap = True
    textbox.fill.solid()
    textbox.fill.fore_color.rgb = RGBColor(240, 240, 255)
    textbox.line.fill.background()

def split_section(title, bullets, max_len=10):
    slides = []
    for i in range(0, len(bullets), max_len):
        part_title = f"{title} (Part {i // max_len + 1})" if i > 0 else title
        slides.append((part_title, bullets[i:i + max_len]))
    return slides

def clean_title(raw_title):
    title = raw_title.strip()
    title = re.sub(r'^(#+|[*_]{2}|Slide\s*\d+:)?\s*', '', title)
    title = title.replace('*', '')
    return title.title() if title else "Untitled Section"

def extract_table(text_block):
    lines = [line.strip() for line in text_block.splitlines() if '|' in line]
    if len(lines) < 2:
        return None
    headers = [flatten_to_string(cell.strip()) for cell in lines[0].split('|') if cell.strip()]
    rows = []
    for line in lines[2:]:
        cells = [flatten_to_string(cell.strip()) for cell in line.split('|') if cell.strip()]
        if len(cells) == len(headers):
            rows.append(cells)
    return [headers] + rows if rows else None

def extract_sections(content):
    sections = re.split(r'\n(?=#+\s|^Slide\s*\d+:|^\*\*.+?\*\*|^\w.+?:)', content, flags=re.MULTILINE)
    structured = []
    for section in sections:
        section = section.strip()
        if not section:
            continue
        lines = section.splitlines()
        if len(lines) < 2:
            continue
        title = clean_title(lines[0])
        section_body = "\n".join(lines[1:])
        table = extract_table(section_body)
        if table:
            structured.append(("TABLE", title, table))
        else:
            bullets = []
            for line in lines[1:]:
                clean_line = re.sub(r'^[-•*0-9. ]+', '', flatten_to_string(line)).replace('*', '')
                if clean_line:
                    bullets.extend(textwrap.wrap(clean_line, width=90))
            structured.extend([("TEXT", t, b) for t, b in split_section(title, bullets)])
    return structured

def set_slide_background(slide, color=RGBColor(255, 255, 255)):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title(slide, title_text):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, TITLE_TOP, SLIDE_WIDTH, TITLE_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_BG_COLOR
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = flatten_to_string(title_text).replace('*', '')
    p.font.size = TITLE_FONT_SIZE
    p.font.name = "Segoe UI"
    p.font.bold = True
    p.font.color.rgb = TITLE_FONT_COLOR
    p.alignment = PP_ALIGN.CENTER

def add_bullets(slide, bullet_lines):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CONTENT_BG_COLOR
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(bullet_lines):
        s_val = flatten_to_string(line)
        p = tf.add_paragraph()
        p.text = s_val.strip().replace('*', '')
        p.font.name = "Segoe UI"
        p.font.size = BULLET_FONT_SIZE
        p.font.color.rgb = TEXT_COLOR
        p.margin_left = BULLET_INDENT
        p.font.bold = s_val.strip().endswith(":")
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(2 if i > 0 else 0)
        p.space_after = Pt(6)

def add_table(slide, table_data):
    rows = len(table_data)
    cols = len(table_data[0])
    table_shape = slide.shapes.add_table(rows, cols, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT)
    table = table_shape.table
    for col, value in enumerate(table_data[0]):
        s_val = flatten_to_string(value)
        cell = table.cell(0, col)
        cell.text = s_val.replace('*', '')
        para = cell.text_frame.paragraphs[0]
        para.font.size = TABLE_FONT_SIZE
        para.font.bold = True
        para.font.name = "Segoe UI"
        para.font.color.rgb = TABLE_HEADER_FONT_COLOR
        para.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
    for row_idx, row in enumerate(table_data[1:], start=1):
        for col_idx, value in enumerate(row):
            s_val = flatten_to_string(value)
            cell = table.cell(row_idx, col_idx)
            cell.text = s_val.replace('*', '')
            para = cell.text_frame.paragraphs[0]
            para.font.name = "Segoe UI"
            para.font.size = TABLE_FONT_SIZE
            para.alignment = PP_ALIGN.LEFT
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_ALT_BG

def add_references_slide(prs, references, page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, CONTENT_BG_COLOR)
    add_title(slide, "References")
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CONTENT_BG_COLOR
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    for ref in references:
        s_val = flatten_to_string(ref)
        p = tf.add_paragraph()
        p.text = s_val.replace('*', '')
        p.font.name = "Segoe UI"
        p.font.size = BULLET_FONT_SIZE
        p.font.color.rgb = TEXT_COLOR
        p.alignment = PP_ALIGN.LEFT
        p.margin_left = BULLET_INDENT
        p.space_after = Pt(4)
    add_footer(slide, page_num)

def generate_ppt(content, output_path="outputs/output.pptx", references=None, filename_title="Untitled Document"):
    prs = Presentation()
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    add_cover_slide(prs, filename_title)
    page_num = 2
    sections = extract_sections(content)
    if not sections:
        print("❌ Error: No valid sections found.")
        return
    layout = prs.slide_layouts[6]
    for kind, title, body in sections:
        slide = prs.slides.add_slide(layout)
        set_slide_background(slide)
        add_title(slide, title)
        if kind == "TEXT":
            add_bullets(slide, body)
        elif kind == "TABLE":
            add_table(slide, body)
        add_footer(slide, page_num)
        page_num += 1
    if references:
        add_references_slide(prs, references, page_num)
    prs.save(output_path)
    print(f"✅ PPT saved to: {output_path}")
